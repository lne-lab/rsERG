import os
import pickle
import numpy as np
import matplotlib.pyplot as plt
import math

import ipywidgets as widgets
from ipyfilechooser import FileChooser
from IPython.display import display, clear_output

# Check for python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed, PPT export will be disabled.")

###############################################################################
# 1) Helper Functions
###############################################################################

def exclude_traces(
    psd_array,
    freqs,
    low_band=(1, 3),
    low_band_threshold=3.0,
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,
    test_band_count_threshold=None
):
    """
    Returns (kept_indices, excluded_indices) after applying
    the specified thresholds for outlier detection.
    """
    if test_band_count_threshold is None:
        test_band_count_threshold = len(test_bands) // 2

    mean_psd = np.mean(psd_array, axis=0)
    excluded_traces = []
    kept_traces = []

    low_band_indices = np.where((freqs >= low_band[0]) & (freqs <= low_band[1]))[0]
    band_indices = [
        np.where((freqs >= band[0]) & (freqs <= band[1]))[0]
        for band in test_bands
    ]

    for i, trace in enumerate(psd_array):
        # 1) Check low-frequency extreme outlier
        if np.any(trace[low_band_indices] > low_band_threshold * mean_psd[low_band_indices]):
            excluded_traces.append(i)
            continue

        # 2) Check repeated suprathreshold events in test bands
        suprathreshold_count = 0
        for indices in band_indices:
            if np.any(trace[indices] > test_band_threshold * mean_psd[indices]):
                suprathreshold_count += 1

        if suprathreshold_count >= test_band_count_threshold:
            excluded_traces.append(i)
        else:
            kept_traces.append(i)

    return kept_traces, excluded_traces


def compute_channel_means(
    psds_dict,
    exclude=False,
    low_band=(1,3),
    low_band_threshold=3.0,
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,
    test_band_count_threshold=None
):
    """
    Compute channel means before and after exclusion.
    Returns a dictionary {channel: (original_mean, new_mean)}.
    """
    channel_means = {}
    for channel, data in psds_dict.items():
        psd = data.get('psd', None)
        freqs = data.get('freqs', None)
        if psd is None or freqs is None:
            print(f"Channel '{channel}' is missing 'psd' or 'freqs' data. Skipping.")
            continue

        original_mean = np.mean(psd, axis=0)
        if exclude:
            kept_traces, _ = exclude_traces(
                psd_array=psd,
                freqs=freqs,
                low_band=low_band,
                low_band_threshold=low_band_threshold,
                test_bands=test_bands,
                test_band_threshold=test_band_threshold,
                test_band_count_threshold=test_band_count_threshold
            )
            if kept_traces:
                new_mean = np.mean(psd[kept_traces], axis=0)
            else:
                new_mean = np.zeros_like(original_mean)
        else:
            new_mean = original_mean.copy()
        
        channel_means[channel] = (original_mean, new_mean)
    
    return channel_means


def compute_group_mean(channel_means, selected_channels):
    """
    Compute the mean of 'new_mean' across selected channels.
    Returns a single array (the mean across those channels).
    """
    means = []
    for ch in selected_channels:
        if ch in channel_means:
            orig_m, new_m = channel_means[ch]
            means.append(new_m)
    if means:
        return np.mean(means, axis=0)
    else:
        # If nothing is selected, return zeros of the correct shape
        if channel_means:
            # Grab shape from first channel
            first_val = next(iter(channel_means.values()))  # (orig, new)
            return np.zeros_like(first_val[0])
        else:
            return None

###############################################################################
# 2) Plotting Utility: Plot each channel with a distinct color
###############################################################################
def plot_individual_channels(
    ax,
    channels,
    channel_means_dict,
    freqs,
    colors_for_channels,
    show_original_mean=True,
    show_new_mean=True,
    title="",
    axis_fs=10,
    legend_fs=8,
    tick_fs=8
):
    """
    Plots each channel in the specified list on the given Axes 'ax',
    using a distinct color from 'colors_for_channels'.
    If show_original_mean=True, it plots original mean in that color (solid).
    If show_new_mean=True, it plots new/excluded mean in that color (dashed).
    """
    for ch in channels:
        if ch not in channel_means_dict:
            continue
        orig_mean, new_mean = channel_means_dict[ch]
        color = colors_for_channels.get(ch, "blue")  # fallback if missing

        lbl_orig = f"{ch} Orig"
        lbl_new  = f"{ch} New"

        # Plot original
        if show_original_mean:
            ax.plot(freqs, orig_mean, color=color, linestyle='-', label=lbl_orig)
        # Plot new
        if show_new_mean:
            ax.plot(freqs, new_mean, color=color, linestyle='--', label=lbl_new)
    
    ax.set_title(title, fontsize=axis_fs)
    ax.set_xlabel("Frequency (Hz)", fontsize=axis_fs)
    ax.set_ylabel("PSD (V²/Hz)", fontsize=axis_fs)
    ax.tick_params(axis='both', labelsize=tick_fs)
    ax.legend(fontsize=legend_fs, loc="upper right")


###############################################################################
# 3) Main Interactive GUI
###############################################################################
def build_exportable_plot_psd_gui():
    """
    An interactive GUI to:
      - Load PSD data (via ipyfilechooser)
      - Select channels
      - Show/hide Eye1/Eye2
      - Show/hide original vs new means
      - Plot them in separate Eye1/Eye2 subplots
      - Compute & display final group means, then a separate figure with ONLY 
        the Eye1/Eye2 final means (no individual channels).
      - Export the subplot figures (via ipyfilechooser-based path)
      - Export the final "mean of means" data/figure (also via filechooser)
    """
    # Define these variables in the outer function scope so we can use "nonlocal" later:
    final_data_dict = {}  # will store the final means arrays
    final_fig = None      # will store the final "means only" figure
    current_figures = []  # store the subplots from the "screening" and "final" plots

    # 1. Load PSD Pickle (with ipyfilechooser)
    load_psd_button = widgets.Button(
        description='Load PSD Pickle',
        button_style='info',
        tooltip='Load a pickled PSD file'
    )
    psd_file_chooser = FileChooser(
        os.getcwd(),
        title='Select PSD Pickle File',
        select_default=False
    )
    psd_file_chooser.show_only_files = True
    psd_file_chooser.filter_pattern = ['*.pkl']

    # 2. Select Channels
    channels_dropdown = widgets.SelectMultiple(
        options=[],
        description='Select Channels:',
        layout=widgets.Layout(width='300px', height='200px')
    )

    # 3. Plot toggles
    show_eye1_cb = widgets.Checkbox(value=True, description='Show Eye1 (Ch1-8)')
    show_eye2_cb = widgets.Checkbox(value=True, description='Show Eye2 (Ch9-16)')
    show_original_mean_cb = widgets.Checkbox(value=True, description='Show Original Means')
    show_new_mean_cb = widgets.Checkbox(value=True, description='Show New Means')
    show_kept_cb = widgets.Checkbox(value=True, description='[Unused in this example] Show Kept Traces')
    show_excluded_cb = widgets.Checkbox(value=True, description='[Unused in this example] Show Excluded Traces')
    show_vertical_lines_cb = widgets.Checkbox(value=True, description='[Unused] Show Vertical Lines')

    # 4. Vertical lines (placeholder if needed)
    vertical_lines_text = widgets.Text(
        value="10,15",
        description='Vertical Lines (Hz):',
        layout=widgets.Layout(width='300px')
    )

    # 5. Axis Ranges
    x_min_widget = widgets.FloatText(value=None, description='X-axis Min:', layout=widgets.Layout(width='150px'))
    x_max_widget = widgets.FloatText(value=None, description='X-axis Max:', layout=widgets.Layout(width='150px'))
    y_min_widget = widgets.FloatText(value=None, description='Y-axis Min:', layout=widgets.Layout(width='150px'))
    y_max_widget = widgets.FloatText(value=None, description='Y-axis Max:', layout=widgets.Layout(width='150px'))

    # 6. Font sizes
    title_fs_widget = widgets.IntText(value=10, description='Title FS:', layout=widgets.Layout(width='120px'))
    axis_fs_widget = widgets.IntText(value=8, description='Axis FS:', layout=widgets.Layout(width='120px'))
    legend_fs_widget = widgets.IntText(value=8, description='Legend FS:', layout=widgets.Layout(width='120px'))
    tick_fs_widget = widgets.IntText(value=8, description='Tick FS:', layout=widgets.Layout(width='120px'))
    max_title_len_widget = widgets.IntText(value=40, description='Max Title:', layout=widgets.Layout(width='120px'))

    # 7. Colors (for kept/excluded) - though not heavily used here
    color_kept_widget = widgets.ColorPicker(value='lightgray', description='Kept Trace:', layout=widgets.Layout(width='150px'))
    color_excl_widget = widgets.ColorPicker(value='red', description='Excluded Trace:', layout=widgets.Layout(width='150px'))
    color_oldmean_widget = widgets.ColorPicker(value='blue', description='Original Mean:', layout=widgets.Layout(width='150px'))
    color_newmean_widget = widgets.ColorPicker(value='green', description='New Mean:', layout=widgets.Layout(width='150px'))

    # 8. Test band text
    test_band_text = widgets.Textarea(
        value="(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)",
        description='Test Bands:',
        layout=widgets.Layout(width='300px', height='80px')
    )

    # 9. Number columns
    num_cols_widget = widgets.IntText(value=2, description='Cols/Row:', layout=widgets.Layout(width='130px'))

    # 10. Vertical line alpha (unused in this example)
    vertical_line_alpha_widget = widgets.FloatSlider(
        value=0.6, min=0.0, max=1.0, step=0.1, description='Vertical Line Alpha:',
        continuous_update=False, layout=widgets.Layout(width='300px')
    )

    # 11. Plot Button
    plot_psd_button = widgets.Button(
        description='Plot PSDs',
        button_style='success',
        tooltip='Plot the selected PSDs'
    )

    # 12. Export controls (subplots) via FileChooser
    export_fig_chooser = FileChooser(
        os.getcwd(),
        title='Select or type a file name to export subplots (e.g. MyPlots.png or MySlides.pptx)',
        select_default=False
    )
    # Example: export_fig_chooser.filter_pattern = ['*.png','*.pptx','*.svg','*.jpeg']

    export_button = widgets.Button(
        description='Export Figures',
        button_style='warning',
        tooltip='Export the plotted subplots'
    )

    # 13. Output areas
    load_output_area = widgets.Output()
    plot_output_area = widgets.Output()
    export_output_area = widgets.Output()

    # 14. Exclusion param widgets
    low_band_threshold_widget = widgets.FloatText(
        value=3.0, description='Low Band Threshold:', layout=widgets.Layout(width='150px')
    )
    test_band_threshold_widget = widgets.FloatText(
        value=10.0, description='Test Band Threshold:', layout=widgets.Layout(width='150px')
    )

    # This dictionary will hold the loaded PSD data from the pickle
    loaded_psd = {}

    ########################################################################
    # A) LOAD PSD PICKLE CALLBACK
    ########################################################################
    def on_load_psd_clicked(b):
        with load_output_area:
            clear_output()
            psd_path = psd_file_chooser.selected
            if not psd_path:
                print("Please select a PSD pickle file.")
                return
            if not os.path.isfile(psd_path):
                print(f"The file does not exist: {psd_path}")
                return
            try:
                with open(psd_path, 'rb') as f:
                    loaded_psd.clear()
                    loaded_psd.update(pickle.load(f))
                print(f"Successfully loaded PSD data from '{psd_path}'.")

                # Populate channels dropdown
                # Attempt to sort by numeric channel name if they are like "Ch1", "Ch2", ...
                # Otherwise, just do a normal alphabetical sort
                def ch_sort_key(ch):
                    return int(ch.replace("Ch", "")) if ch.startswith("Ch") else 999999
                channels = sorted(loaded_psd.keys(), key=ch_sort_key)
                
                if not channels:
                    print("The loaded PSD data is empty.")
                    return
                channels_dropdown.options = channels
                print(f"Available Channels: {channels}")
            except Exception as e:
                print(f"Failed to load PSD pickle: {e}")

    load_psd_button.on_click(on_load_psd_clicked)

    ########################################################################
    # B) PLOT PSD BUTTON CALLBACK
    ########################################################################
    def on_plot_psd_clicked(b):
        nonlocal current_figures, final_data_dict, final_fig
        with plot_output_area:
            clear_output()
            current_figures.clear()
            final_data_dict.clear()
            final_fig = None

            selected_channels = list(channels_dropdown.value)
            if not selected_channels:
                print("Please select at least one channel.")
                return
            
            # Toggles
            show_eye1 = show_eye1_cb.value
            show_eye2 = show_eye2_cb.value
            show_orig = show_original_mean_cb.value
            show_new  = show_new_mean_cb.value

            # Axis ranges
            x_min = x_min_widget.value
            x_max = x_max_widget.value
            y_min = y_min_widget.value
            y_max = y_max_widget.value

            # Font sizes
            title_fs = title_fs_widget.value
            axis_fs  = axis_fs_widget.value
            legend_fs= legend_fs_widget.value
            tick_fs  = tick_fs_widget.value
            max_title_len = max_title_len_widget.value

            # Parse test bands
            test_bands_str = test_band_text.value.strip()
            test_bands_list = []
            if test_bands_str:
                test_bands_str = test_bands_str.replace(" ", "")
                pairs = test_bands_str.split(")")
                for p in pairs:
                    p = p.strip(",").strip()
                    p = p.strip("(").strip()
                    if not p:
                        continue
                    vals = p.split(",")
                    if len(vals) == 2:
                        try:
                            lowf = float(vals[0])
                            highf = float(vals[1])
                            test_bands_list.append((lowf, highf))
                        except:
                            pass
                if not test_bands_list:
                    test_bands_list = [(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)]
            else:
                test_bands_list = [(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)]
            
            # Exclusion parameters
            low_band = (1,3)
            low_band_threshold = low_band_threshold_widget.value
            test_band_threshold = test_band_threshold_widget.value

            # Separate Eye1 vs Eye2 channels
            # This logic is user-defined: Eye1=Ch1..Ch8, Eye2=Ch9..Ch16
            def get_ch_num(ch):
                # If the channel is something like "Ch1" => int(...) = 1
                # If not, fallback to 99999
                try:
                    return int(ch.replace("Ch", ""))
                except:
                    return 99999
            eye1_channels = [ch for ch in selected_channels if 1 <= get_ch_num(ch) <= 8]
            eye2_channels = [ch for ch in selected_channels if 9 <= get_ch_num(ch) <= 16]

            # 1) Compute channel means: before & after
            channel_means_before = compute_channel_means(
                psds_dict=loaded_psd,
                exclude=False
            )
            channel_means_after = compute_channel_means(
                psds_dict=loaded_psd,
                exclude=True,
                low_band=low_band,
                low_band_threshold=low_band_threshold,
                test_bands=test_bands_list,
                test_band_threshold=test_band_threshold,
                test_band_count_threshold=None
            )

            # 2) SCREENING STEP
            import matplotlib.cm as cm
            color_cycle = cm.get_cmap('tab20').colors
            colors_for_channels = {}
            for i, ch in enumerate(sorted(selected_channels)):
                colors_for_channels[ch] = color_cycle[i % len(color_cycle)]

            n_subplots = 0
            if show_eye1 and eye1_channels: n_subplots += 1
            if show_eye2 and eye2_channels: n_subplots += 1

            if n_subplots == 0:
                print("No Eye1/Eye2 channels to show in screening step.")
            else:
                fig_screen, axes_screen = plt.subplots(1, n_subplots, figsize=(6*n_subplots, 5))
                if n_subplots == 1:
                    axes_screen = [axes_screen]
                ax_idx = 0

                if show_eye1 and eye1_channels:
                    ax_e1 = axes_screen[ax_idx]
                    freqs_e1 = loaded_psd[eye1_channels[0]]['freqs']
                    plot_individual_channels(
                        ax=ax_e1,
                        channels=eye1_channels,
                        channel_means_dict=channel_means_before,
                        freqs=freqs_e1,
                        colors_for_channels=colors_for_channels,
                        show_original_mean=show_orig,
                        show_new_mean=show_new,
                        title="Screening Eye1",
                        axis_fs=axis_fs,
                        legend_fs=legend_fs,
                        tick_fs=tick_fs
                    )
                    if x_min is not None: ax_e1.set_xlim(left=x_min)
                    if x_max is not None: ax_e1.set_xlim(right=x_max)
                    if y_min is not None: ax_e1.set_ylim(bottom=y_min)
                    if y_max is not None: ax_e1.set_ylim(top=y_max)
                    ax_idx += 1

                if show_eye2 and eye2_channels:
                    ax_e2 = axes_screen[ax_idx]
                    freqs_e2 = loaded_psd[eye2_channels[0]]['freqs']
                    plot_individual_channels(
                        ax=ax_e2,
                        channels=eye2_channels,
                        channel_means_dict=channel_means_before,
                        freqs=freqs_e2,
                        colors_for_channels=colors_for_channels,
                        show_original_mean=show_orig,
                        show_new_mean=show_new,
                        title="Screening Eye2",
                        axis_fs=axis_fs,
                        legend_fs=legend_fs,
                        tick_fs=tick_fs
                    )
                    if x_min is not None: ax_e2.set_xlim(left=x_min)
                    if x_max is not None: ax_e2.set_xlim(right=x_max)
                    if y_min is not None: ax_e2.set_ylim(bottom=y_min)
                    if y_max is not None: ax_e2.set_ylim(top=y_max)

                plt.tight_layout()
                plt.show()
                current_figures.append(fig_screen)

            # 3) Create selection widgets for final group means
            group_selection_widgets = {}
            if eye1_channels:
                group_selection_widgets['Eye1'] = widgets.SelectMultiple(
                    options=eye1_channels,
                    value=eye1_channels,
                    description='Eye1 Channels:',
                    layout=widgets.Layout(width='300px', height='150px')
                )
            if eye2_channels:
                group_selection_widgets['Eye2'] = widgets.SelectMultiple(
                    options=eye2_channels,
                    value=eye2_channels,
                    description='Eye2 Channels:',
                    layout=widgets.Layout(width='300px', height='150px')
                )

            display(widgets.HTML("<h3>Select Channels to Include in Final Group Mean(s)</h3>"))
            group_sel_box = widgets.VBox(list(group_selection_widgets.values()))
            display(group_sel_box)

            final_plot_button = widgets.Button(
                description='Plot Final Group Means',
                button_style='success'
            )
            display(final_plot_button)

            out_final_plots = widgets.Output()
            display(out_final_plots)

            def on_final_plot_button_clicked(btn):
                nonlocal final_data_dict, final_fig, current_figures
                with out_final_plots:
                    clear_output()

                    final_eye1 = []
                    final_eye2 = []
                    if 'Eye1' in group_selection_widgets:
                        final_eye1 = list(group_selection_widgets['Eye1'].value)
                    if 'Eye2' in group_selection_widgets:
                        final_eye2 = list(group_selection_widgets['Eye2'].value)

                    group_mean_eye1 = compute_group_mean(channel_means_after, final_eye1) if final_eye1 else None
                    group_mean_eye2 = compute_group_mean(channel_means_after, final_eye2) if final_eye2 else None

                    # 4) Plot final channels in separate Eye1/Eye2 subplots
                    n_subplots_final = 0
                    if show_eye1 and final_eye1: n_subplots_final += 1
                    if show_eye2 and final_eye2: n_subplots_final += 1

                    fig_final = None
                    if n_subplots_final == 0:
                        print("No final channels selected or no Eye boxes checked.")
                    else:
                        fig_final, axes_final = plt.subplots(1, n_subplots_final, figsize=(6*n_subplots_final, 5))
                        if n_subplots_final == 1:
                            axes_final = [axes_final]
                        ax_idx_final = 0

                        # Eye1 final
                        if show_eye1 and final_eye1:
                            ax_e1f = axes_final[ax_idx_final]
                            freqs_e1 = loaded_psd[final_eye1[0]]['freqs']
                            import matplotlib.cm as cm
                            color_cycle = cm.get_cmap('tab20').colors
                            colors_for_channels = {}
                            for i, ch in enumerate(sorted(final_eye1)):
                                colors_for_channels[ch] = color_cycle[i % len(color_cycle)]
                            plot_individual_channels(
                                ax=ax_e1f,
                                channels=final_eye1,
                                channel_means_dict=channel_means_after,
                                freqs=freqs_e1,
                                colors_for_channels=colors_for_channels,
                                show_original_mean=False,
                                show_new_mean=True,
                                title="Final Eye1",
                                axis_fs=axis_fs,
                                legend_fs=legend_fs,
                                tick_fs=tick_fs
                            )
                            if group_mean_eye1 is not None:
                                ax_e1f.plot(freqs_e1, group_mean_eye1, color='black', linewidth=2, label='Eye1 Group Mean')
                                ax_e1f.legend(fontsize=legend_fs, loc="upper right")

                            if x_min is not None: ax_e1f.set_xlim(left=x_min)
                            if x_max is not None: ax_e1f.set_xlim(right=x_max)
                            if y_min is not None: ax_e1f.set_ylim(bottom=y_min)
                            if y_max is not None: ax_e1f.set_ylim(top=y_max)
                            ax_idx_final += 1

                        # Eye2 final
                        if show_eye2 and final_eye2:
                            ax_e2f = axes_final[ax_idx_final]
                            freqs_e2 = loaded_psd[final_eye2[0]]['freqs']
                            import matplotlib.cm as cm
                            color_cycle = cm.get_cmap('tab20').colors
                            colors_for_channels = {}
                            for i, ch in enumerate(sorted(final_eye2)):
                                colors_for_channels[ch] = color_cycle[i % len(color_cycle)]
                            plot_individual_channels(
                                ax=ax_e2f,
                                channels=final_eye2,
                                channel_means_dict=channel_means_after,
                                freqs=freqs_e2,
                                colors_for_channels=colors_for_channels,
                                show_original_mean=False,
                                show_new_mean=True,
                                title="Final Eye2",
                                axis_fs=axis_fs,
                                legend_fs=legend_fs,
                                tick_fs=tick_fs
                            )
                            if group_mean_eye2 is not None:
                                ax_e2f.plot(freqs_e2, group_mean_eye2, color='black', linewidth=2, label='Eye2 Group Mean')
                                ax_e2f.legend(fontsize=legend_fs, loc="upper right")

                            if x_min is not None: ax_e2f.set_xlim(left=x_min)
                            if x_max is not None: ax_e2f.set_xlim(right=x_max)
                            if y_min is not None: ax_e2f.set_ylim(bottom=y_min)
                            if y_max is not None: ax_e2f.set_ylim(top=y_max)

                        plt.tight_layout()
                        plt.show()
                        if fig_final:
                            current_figures.append(fig_final)

                    # 5) Plot a second "means only" figure: just Eye1/Eye2 final means
                    if (group_mean_eye1 is not None) or (group_mean_eye2 is not None):
                        fig_mom, ax_mom = plt.subplots(figsize=(6,4))

                        # Eye1 final
                        if group_mean_eye1 is not None and len(final_eye1) > 0:
                            freqs_e1 = loaded_psd[final_eye1[0]]['freqs']
                            ax_mom.plot(freqs_e1, group_mean_eye1, color='red', linewidth=2, label='Eye1 Final Mean')
                            # store in final_data_dict
                            final_data_dict["Eye1_MeanOfMeans"] = group_mean_eye1
                            final_data_dict["Eye1_freqs"] = freqs_e1

                        # Eye2 final
                        if group_mean_eye2 is not None and len(final_eye2) > 0:
                            freqs_e2 = loaded_psd[final_eye2[0]]['freqs']
                            ax_mom.plot(freqs_e2, group_mean_eye2, color='blue', linewidth=2, label='Eye2 Final Mean')
                            final_data_dict["Eye2_MeanOfMeans"] = group_mean_eye2
                            final_data_dict["Eye2_freqs"] = freqs_e2

                        ax_mom.set_title("Final Means of Means Only", fontsize=title_fs)
                        ax_mom.set_xlabel("Frequency (Hz)", fontsize=axis_fs)
                        ax_mom.set_ylabel("PSD (V²/Hz)", fontsize=axis_fs)
                        ax_mom.tick_params(axis='both', labelsize=tick_fs)
                        ax_mom.legend(fontsize=legend_fs)
                        if x_min is not None: ax_mom.set_xlim(left=x_min)
                        if x_max is not None: ax_mom.set_xlim(right=x_max)
                        if y_min is not None: ax_mom.set_ylim(bottom=y_min)
                        if y_max is not None: ax_mom.set_ylim(top=y_max)

                        plt.tight_layout()
                        plt.show()
                        current_figures.append(fig_mom)
                        final_fig = fig_mom  # store the final means-only figure

            final_plot_button.on_click(on_final_plot_button_clicked)

    plot_psd_button.on_click(on_plot_psd_clicked)

    ########################################################################
    # C) EXPORT BUTTON CALLBACK (for subplots) - using filechooser
    ########################################################################
    def on_export_button_clicked(b):
        nonlocal current_figures
        with export_output_area:
            clear_output()
            if not current_figures:
                print("No figures to export. Please plot PSDs first.")
                return

            chosen_export_path = export_fig_chooser.selected
            if not chosen_export_path:
                print("No export file chosen. Please select or type a filename (e.g. MyPlot.png or MySlides.pptx).")
                return

            base, ext = os.path.splitext(chosen_export_path)
            ext = ext.lower()
            if not ext:
                print("No file extension detected. Please include something like .png or .pptx.")
                return

            if ext in ('.ppt', '.pptx'):
                if not HAS_PPTX:
                    print("python-pptx not installed. Cannot export PPT.")
                    return
                ppt_file = chosen_export_path
                print(f"Exporting {len(current_figures)} figure(s) to {ppt_file} ...")
                prs = Presentation()
                blank_layout = prs.slide_layouts[6]
                for i, fig in enumerate(current_figures, start=1):
                    temp_png = f"{base}_temp_{i}.png"
                    fig.savefig(temp_png, format='png', dpi=150)
                    slide = prs.slides.add_slide(blank_layout)
                    left = top = Inches(1)
                    slide.shapes.add_picture(temp_png, left, top, width=Inches(8), height=Inches(4.5))
                    os.remove(temp_png)
                prs.save(ppt_file)
                print(f"Done exporting to {ppt_file}")
            else:
                # treat it as an image format
                print(f"Exporting {len(current_figures)} figure(s) as *{ext} images.")
                for i, fig in enumerate(current_figures, start=1):
                    out_file = f"{base}_{i}{ext}"
                    fig.savefig(out_file, format=ext.lstrip('.'), dpi=150)
                    print(f"Saved => {out_file}")
                print("Done exporting images.")

    export_button.on_click(on_export_button_clicked)

    ########################################################################
    # D) Export final "means of means" figure/data (multi-format) - also filechooser
    ########################################################################
    final_export_label = widgets.HTML("<h3>Export Final Means Plot/Data</h3>")

    final_export_chooser = FileChooser(
        os.getcwd(),
        title='Select or type a path for final means export (e.g. FinalData.pkl, .xlsx, .png, .pptx, etc.)',
        select_default=False
    )
    # final_export_chooser.filter_pattern = ['*.pkl','*.xlsx','*.png','*.pptx','*.svg','*.jpeg']

    final_export_button = widgets.Button(description='Export Final Means', button_style='info')
    final_export_output = widgets.Output()

    def on_final_export_clicked(b):
        nonlocal final_data_dict, final_fig
        with final_export_output:
            clear_output()
            chosen_path = final_export_chooser.selected
            if not chosen_path:
                print("No export file chosen. Please select or type a filename (e.g. FinalData.pkl).")
                return

            base, ext = os.path.splitext(chosen_path)
            ext = ext.lower()

            if not ext:
                print("No file extension detected. Please include something like .pkl, .xlsx, .png, or .pptx.")
                return

            if not final_data_dict and final_fig is None:
                print("No final data or final figure found. Please run 'Plot Final Group Means' first.")
                return

            print(f"Exporting final means as {ext} ...")

            # 1) If extension = .pkl => pickle final_data_dict
            if ext == '.pkl':
                try:
                    with open(chosen_path, 'wb') as f:
                        pickle.dump(final_data_dict, f)
                    print(f"Saved final_data_dict => {chosen_path}")
                except Exception as e:
                    print(f"ERROR saving pickle: {e}")

            # 2) If extension = .xlsx => export Excel
            elif ext == '.xlsx' or ext == '.xls':
                try:
                    import pandas as pd
                    df = pd.DataFrame()
                    for k,v in final_data_dict.items():
                        arr = np.array(v)
                        df[k] = arr.ravel()  # flatten to 1D
                    df.to_excel(chosen_path, index=False)
                    print(f"Saved final_data_dict => {chosen_path}")
                except Exception as e:
                    print(f"ERROR saving Excel: {e}")

            # 3) If extension in image formats => save the final_fig
            elif ext in ('.png','.svg','.jpeg','.jpg'):
                if final_fig is None:
                    print("No final figure to export.")
                    return
                try:
                    final_fig.savefig(chosen_path, format=ext.lstrip('.'), dpi=150)
                    print(f"Saved final_fig => {chosen_path}")
                except Exception as e:
                    print(f"ERROR saving figure: {e}")

            # 4) If extension = .ppt or .pptx => PPT
            elif ext in ('.ppt','.pptx'):
                if not HAS_PPTX:
                    print("python-pptx not installed. Cannot export PPT.")
                    return
                if final_fig is None:
                    print("No final figure to export.")
                    return
                try:
                    prs = Presentation()
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    tmp_png = f"{base}_temp.png"
                    final_fig.savefig(tmp_png, dpi=150)
                    left = top = Inches(1)
                    slide.shapes.add_picture(tmp_png, left, top, width=Inches(8), height=Inches(4.5))
                    os.remove(tmp_png)
                    prs.save(chosen_path)
                    print(f"Exported final_fig => {chosen_path}")
                except Exception as e:
                    print(f"ERROR exporting PPT: {e}")

            else:
                print(f"Unsupported file extension '{ext}'. Please use .pkl, .xlsx, .png, .ppt, .pptx, .svg, or .jpeg")

            print("Done exporting final means data/plot.")

    final_export_button.on_click(on_final_export_clicked)

    ########################################################################
    # E) LAYOUT & DISPLAY
    ########################################################################
    exclusion_params_box = widgets.HBox([low_band_threshold_widget, test_band_threshold_widget])
    test_bands_box = widgets.VBox([
        widgets.Label("Test Band Ranges: (format: (7,9),(9,11), etc.)"),
        test_band_text
    ])
    axis_range_box = widgets.HBox([x_min_widget, x_max_widget, y_min_widget, y_max_widget])
    font_size_box = widgets.HBox([title_fs_widget, axis_fs_widget, legend_fs_widget, tick_fs_widget, max_title_len_widget])
    color_box = widgets.HBox([color_kept_widget, color_excl_widget, color_oldmean_widget, color_newmean_widget])
    vertical_line_alpha_box = widgets.HBox([vertical_line_alpha_widget])

    plot_controls_box = widgets.VBox([
        widgets.HTML("<b>Plotting Options:</b>"),
        show_eye1_cb,
        show_eye2_cb,
        show_original_mean_cb,
        show_new_mean_cb,
        show_kept_cb,
        show_excluded_cb,
        show_vertical_lines_cb,
        vertical_lines_text,
        exclusion_params_box,
        vertical_line_alpha_box,
        test_bands_box,
        axis_range_box,
        font_size_box,
        color_box,
        widgets.HBox([num_cols_widget])
    ])

    export_controls = widgets.VBox([
        widgets.HTML("<h3>Export Subplots (Screening/Final) Figures</h3>"),
        export_fig_chooser,
        export_button,
        export_output_area,

        widgets.HTML("<h3>Export Final Means Plot/Data</h3>"),
        final_export_label,
        final_export_chooser,
        final_export_button,
        final_export_output
    ])

    ui = widgets.VBox([
        widgets.HTML("<h2>Enhanced PSD Plotter GUI with ipyfilechooser Exports</h2>"),
        
        widgets.HTML("<h3>1) Load PSD Pickle</h3>"),
        widgets.HBox([load_psd_button, psd_file_chooser]),
        load_output_area,
        
        widgets.HTML("<h3>2) Select Channels</h3>"),
        channels_dropdown,
        
        widgets.HTML("<h3>3) Configure Plot Options & Plot</h3>"),
        plot_controls_box,
        plot_psd_button,
        plot_output_area,
        
        widgets.HTML("<h3>4) Exporting Subplots & Final Means</h3>"),
        export_controls
    ])

    display(ui)

# Build and display the GUI
build_exportable_plot_psd_gui()
