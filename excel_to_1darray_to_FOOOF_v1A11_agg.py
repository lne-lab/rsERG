import os
import re
import pickle
import pandas as pd
import numpy as np
import ipywidgets as widgets
from IPython.display import display, clear_output
import matplotlib.pyplot as plt
from fooof import FOOOFGroup
from scipy.spatial.distance import cdist
from tqdm.notebook import tqdm

# Tkinter for file/directory dialogs
from tkinter import Tk, filedialog

# Optional: For PPT export
try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

###############################################################################
#                           CORE ANALYSIS FUNCTIONS
###############################################################################

def process_single_sheet_excel(file_paths):
    """
    For *grouped* data:
    Reads each .xlsx file (assuming each file has exactly ONE relevant sheet).
      - The first column is assumed to be Frequency.
      - Each subsequent column is a PSD for a particular group or mouse.

    Returns a dict of the form:
      {
         file_path: { 
            "freq": freq_array,
            "Mouse1_Eye1": PSD_array,
            "Mouse1_Eye2": PSD_array,
            ...
         }
      }
    """
    psd_data_dict_all = {}

    for file_path in file_paths:
        try:
            # Assume the Excel has only ONE relevant sheet; read the first sheet by default
            df = pd.read_excel(file_path, sheet_name=0)

            # The first column is frequency
            freq_column = df.iloc[:, 0].values

            # Prepare a dictionary for all PSD columns
            psd_data_dict = {"freq": freq_column}

            # For each remaining column, store it as a separate PSD vector
            for col_name in df.columns[1:]:
                psd_data_dict[col_name] = df[col_name].values

            psd_data_dict_all[file_path] = psd_data_dict

        except Exception as e:
            print(f"An error occurred while processing {file_path}: {e}")

    return psd_data_dict_all

def process_eye_averages_excel(file_paths):
    """
    For *individual* mice (old approach):
    Reads the sheet named "Eye Averages" (the 17th sheet, but we reference by name).
    Expects columns:
      - Column A: Frequencies
      - Column B: "Eye1 Average PSD"
      - Column C: "Eye2 Average PSD"

    Returns a dict of the form:
      {
         file_path: {
            "freq": freq_array,
            "Eye1 Average PSD": psd_array_for_eye1,
            "Eye2 Average PSD": psd_array_for_eye2
         }
      }
    """
    psd_data_dict_all = {}

    for file_path in file_paths:
        try:
            df = pd.read_excel(file_path, sheet_name="Eye Averages")

            freq_column = df.iloc[:, 0].values
            eye1_col = df["Eye1 Average PSD"].values
            eye2_col = df["Eye2 Average PSD"].values

            psd_data_dict = {
                "freq": freq_column,
                "Eye1 Average PSD": eye1_col,
                "Eye2 Average PSD": eye2_col
            }

            psd_data_dict_all[file_path] = psd_data_dict

        except Exception as e:
            print(f"An error occurred while processing 'Eye Averages' in {file_path}: {e}")

    return psd_data_dict_all

def process_individual_mice_means_excel(file_paths):
    """
    For *individual mice means* (NEW approach):
    Reads exactly one sheet, expecting four columns:
      - Column A: PSD for Eye1
      - Column B: Frequency (for Eye1)
      - Column C: PSD for Eye2
      - Column D: Frequency (redundant, for Eye2)

    We use Column A & C for PSD, and Column B for Frequency.
    Column D is ignored since it's redundant.

    Returns a dict of the form:
      {
         file_path: {
            "freq": freq_array,
            "Eye1": psd_array_for_eye1,
            "Eye2": psd_array_for_eye2
         }
      }
    """
    psd_data_dict_all = {}

    for file_path in file_paths:
        try:
            df = pd.read_excel(file_path, sheet_name=0)
            
            # Column A: PSD for Eye1
            eye1_col = df.iloc[:, 0].values
            # Column B: frequency
            freq_col = df.iloc[:, 1].values
            # Column C: PSD for Eye2
            eye2_col = df.iloc[:, 2].values
            # Column D: freq again (redundant), so ignore

            psd_data_dict = {
                "freq": freq_col,
                "Eye1": eye1_col,
                "Eye2": eye2_col
            }

            psd_data_dict_all[file_path] = psd_data_dict

        except Exception as e:
            print(f"An error occurred while processing {file_path}: {e}")

    return psd_data_dict_all

def run_fooof_analysis(psd_data_dict_all, freq_range, amp_threshold, r2_threshold, max_peaks,
                       fitting_mode, peak_width_limits):
    """
    Runs FOOOF analysis on each PSD column (except for 'freq') in each file.
    Returns a dictionary (fg_dict) mapping a unique key -> FOOOFGroup object.
    The unique key is constructed as: "columnName_fileName".
    """
    fg_dict = {}

    for file_path, psd_data_dict in tqdm(psd_data_dict_all.items(), desc="Fitting PSDs with FOOOF"):
        freq = psd_data_dict["freq"]
        file_suffix = os.path.basename(file_path).replace(".xlsx", "")

        # For each PSD column in the dictionary
        for col_name, psd in psd_data_dict.items():
            if col_name == "freq":
                continue  # skip the frequency array

            fg = FOOOFGroup(
                peak_width_limits=peak_width_limits,
                max_n_peaks=max_peaks,
                min_peak_height=amp_threshold,
                verbose=False,
                aperiodic_mode=fitting_mode,
            )

            # FOOOF expects a 2D array if fitting multiple spectra
            # Here, we pass only 1 PSD -> shape (1, n_freqs)
            fg.fit(freq, psd[np.newaxis, :], freq_range)

            # Drop any fits below the R² threshold
            fg.drop(fg.get_params('r_squared') < r2_threshold)

            dict_key = f"{col_name}_{file_suffix}"
            fg_dict[dict_key] = fg

    return fg_dict

def plot_closest_to_mean(fg, sheet_name, export_dir, base_filename, formats,
                         x_min, x_max, y_min, y_max,
                         show_grid=True, x_tick_font_size=8,
                         include_r2=False, include_peak_table=False):
    """
    For FOOOFGroup objects with multiple PSD fits, selects 10 spectra closest
    to the mean (in Euclidean sense) and plots them.
    """
    models = [fg.get_fooof(ind=i) for i in range(len(fg))]
    spectra = [model.power_spectrum for model in models]
    spectra = np.array(spectra)

    if spectra.ndim != 2:
        raise ValueError(f"Unexpected spectra shape: {spectra.shape}. Expected 2D array.")

    mean_spectrum = np.mean(spectra, axis=0)
    distances = cdist(spectra, mean_spectrum[None, :], metric='euclidean').flatten()
    closest_indices = np.argsort(distances)[:10]

    # If exporting to PPT:
    if "ppt" in formats and PPTX_AVAILABLE:
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # blank

    for i, index in enumerate(closest_indices):
        fm = fg.get_fooof(ind=index, regenerate=True)
        r2_value = fm.get_params('r_squared')
        if include_r2:
            title = f"{sheet_name} - Closest Example {index} (R²: {r2_value:.2f})"
        else:
            title = f"{sheet_name} - Closest Example {index}"

        fm.plot(title=title, plot_peaks='shade')
        plt.xlim(x_min, x_max)
        plt.ylim(y_min, y_max)
        x_min_int = int(np.floor(x_min))
        x_max_int = int(np.ceil(x_max))
        plt.xticks(range(x_min_int, x_max_int + 1, 1))
        plt.setp(plt.gca().get_xticklabels(), fontsize=x_tick_font_size)
        plt.grid(show_grid)
        plt.margins(0)

        fig = plt.gcf()
        file_tag = f"{base_filename}_{sheet_name}_ex{i}"

        # Save or export
        for fmt in formats:
            if fmt in ["png", "svg", "jpeg"]:
                filename = os.path.join(export_dir, f"{file_tag}.{fmt}")
                fig.savefig(filename, dpi=300, bbox_inches='tight')
            elif fmt == "ppt" and PPTX_AVAILABLE:
                tmp_png = os.path.join(export_dir, f"{file_tag}_ppt_temp.png")
                fig.savefig(tmp_png, dpi=300, bbox_inches='tight')
                slide = prs.slides.add_slide(slide_layout)
                left = top = Inches(1)
                slide.shapes.add_picture(tmp_png, left, top, height=Inches(5))

        plt.show()

        # Optionally display peak parameters as a table
        if include_peak_table:
            peak_params = fm.peak_params_
            if peak_params.size > 0:
                df_peaks = pd.DataFrame(peak_params, columns=['Center Freq (Hz)', 'Amplitude', 'FWHM'])
                print("Detected Peaks:")
                display(df_peaks)
            else:
                print("No peaks detected.")

        plt.clf()

    # Save PPT if requested
    if "ppt" in formats and PPTX_AVAILABLE:
        pptx_filename = os.path.join(export_dir, f"{base_filename}_{sheet_name}.pptx")
        prs.save(pptx_filename)

def plot_all_psds(fg, sheet_name, export_dir, base_filename, formats,
                  x_min, x_max, y_min, y_max,
                  show_grid=True, x_tick_font_size=8,
                  include_r2=False, include_peak_table=False):
    """
    Plots each fitted PSD in the FOOOFGroup. Typically each FOOOFGroup here has 1 PSD,
    but if there's more, it plots each.
    """
    n_fits = len(fg)
    if "ppt" in formats and PPTX_AVAILABLE:
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]

    for i in range(n_fits):
        fm = fg.get_fooof(ind=i, regenerate=True)
        r2_value = fm.get_params('r_squared')
        if include_r2:
            title = f"{sheet_name} - Example {i} (R²: {r2_value:.2f})"
        else:
            title = f"{sheet_name} - Example {i}"

        fm.plot(title=title, plot_peaks='shade')
        plt.xlim(x_min, x_max)
        plt.ylim(y_min, y_max)
        x_min_int = int(np.floor(x_min))
        x_max_int = int(np.ceil(x_max))
        plt.xticks(range(x_min_int, x_max_int+1, 1))
        plt.setp(plt.gca().get_xticklabels(), fontsize=x_tick_font_size)
        plt.grid(show_grid)
        plt.margins(0)

        fig = plt.gcf()
        file_tag = f"{base_filename}_{sheet_name}_ex{i}"
        for fmt in formats:
            if fmt in ["png", "svg", "jpeg"]:
                filename = os.path.join(export_dir, f"{file_tag}.{fmt}")
                fig.savefig(filename, dpi=300, bbox_inches='tight')
            elif fmt == "ppt" and PPTX_AVAILABLE:
                tmp_png = os.path.join(export_dir, f"{file_tag}_ppt_temp.png")
                fig.savefig(tmp_png, dpi=300, bbox_inches='tight')
                slide = prs.slides.add_slide(slide_layout)
                left = top = Inches(1)
                slide.shapes.add_picture(tmp_png, left, top, height=Inches(5))
        plt.show()

        # Optionally display peak parameters
        if include_peak_table:
            peak_params = fm.peak_params_
            if peak_params.size > 0:
                df_peaks = pd.DataFrame(peak_params, columns=['Center Freq (Hz)', 'Amplitude', 'FWHM'])
                print("Detected Peaks:")
                display(df_peaks)
            else:
                print("No peaks detected.")
        plt.clf()

    if "ppt" in formats and PPTX_AVAILABLE:
        pptx_filename = os.path.join(export_dir, f"{base_filename}_{sheet_name}.pptx")
        prs.save(pptx_filename)

def save_fooof_group(fg_dict, filename="fooof_groups.pkl"):
    """
    Serializes and saves the fg_dict (mapping keys -> FOOOFGroup) to a pickle file.
    """
    with open(filename, "wb") as f:
        pickle.dump(fg_dict, f)
    print(f"FOOOFGroup dictionary saved to {filename}")

def export_fg_dict_excel(fg_dict, excel_filename):
    """
    Exports FOOOFGroup parameters to an Excel file.
    """
    rows = []
    for key, fg in fg_dict.items():
        n_fits = len(fg)
        r2_params = fg.get_params('r_squared')
        ap_params = fg.get_params('aperiodic_params')
        peak_params = fg.get_params('peak_params')
        
        for i in range(n_fits):
            r2 = r2_params[i] if i < len(r2_params) else None
            ap = ap_params[i] if i < len(ap_params) else None
            peaks = peak_params[i] if i < len(peak_params) else None
            
            row = {
                "dict_key": key,
                "fit_index": i,
                "r_squared": r2,
                "aperiodic_params": str(ap),
                "peak_params": str(peaks)
            }
            rows.append(row)
    df = pd.DataFrame(rows)
    df.to_excel(excel_filename, index=False)
    print(f"FOOOF results exported to Excel: {excel_filename}")

def display_fg_dict_table(fg_dict):
    """
    Displays FOOOFGroup parameters in a table within the notebook.
    """
    rows = []
    for key, fg in fg_dict.items():
        n_fits = len(fg)
        r2_params = fg.get_params('r_squared')
        ap_params = fg.get_params('aperiodic_params')
        peak_params = fg.get_params('peak_params')
        for i in range(n_fits):
            r2 = r2_params[i] if i < len(r2_params) else None
            ap = ap_params[i] if i < len(ap_params) else None
            peaks = peak_params[i] if i < len(peak_params) else None
            row = {
                "dict_key": key,
                "fit_index": i,
                "r_squared": r2,
                "aperiodic_params": str(ap),
                "peak_params": str(peaks)
            }
            rows.append(row)
    df = pd.DataFrame(rows)
    print("FOOOF Group Export Table:")
    display(df)

###############################################################################
#                         TKINTER-BASED GUI CODE & WIDGETS
###############################################################################

selected_files = []
export_directory = ""
fg_dict = {}  # Will store final FOOOFGroup results

###############################################################################
# Widgets
###############################################################################

# 1) Analysis mode: Now has a third option 'Individual Mice Means'
analysis_mode = widgets.Dropdown(
    options=["Grouped Mice", "Individual Mice", "Individual Mice Means"],
    value="Grouped Mice",
    description="Analysis Mode:",
    layout=widgets.Layout(width="50%")
)

# Parameter widgets
freq_range_min = widgets.FloatText(description="Freq Min:", value=4.0, layout=widgets.Layout(width="30%"))
freq_range_max = widgets.FloatText(description="Freq Max:", value=45.0, layout=widgets.Layout(width="30%"))
amp_threshold = widgets.FloatText(description="Amplitude Threshold:", value=0.2, layout=widgets.Layout(width="50%"))
r2_threshold = widgets.FloatText(description="R² Threshold:", value=0.5, layout=widgets.Layout(width="50%"))
max_peaks = widgets.IntText(description="Max Peaks:", value=2, layout=widgets.Layout(width="50%"))
fitting_mode = widgets.Dropdown(description="Fitting Mode:", options=["knee", "fixed"], value="knee", layout=widgets.Layout(width="50%"))
peak_width_min = widgets.FloatText(description="Peak Width Min:", value=2.0, layout=widgets.Layout(width="50%"))
peak_width_max = widgets.FloatText(description="Peak Width Max:", value=10.0, layout=widgets.Layout(width="50%"))

# Axis-range widgets
x_axis_min = widgets.FloatText(description="X Axis Min:", value=4.0, layout=widgets.Layout(width="33%"))
x_axis_max = widgets.FloatText(description="X Axis Max:", value=45.0, layout=widgets.Layout(width="33%"))
y_axis_min = widgets.FloatText(description="Y Axis Min:", value=0.0, layout=widgets.Layout(width="33%"))
y_axis_max = widgets.FloatText(description="Y Axis Max:", value=10.0, layout=widgets.Layout(width="33%"))

# Show Grid
grid_checkbox = widgets.Checkbox(value=True, description="Show Grid")

# X-Tick Font Size
x_tick_font_size_widget = widgets.IntSlider(
    value=8,
    min=6,
    max=20,
    step=1,
    description="X Tick Font Size:",
    readout=True,
    layout=widgets.Layout(width="50%")
)

# Include R² in plot title
include_r2_checkbox = widgets.Checkbox(value=True, description="Include R² in Plot Titles")

# Include Peak Table
include_peak_table_checkbox = widgets.Checkbox(value=True, description="Include Peak Params Table")

# Export figures?
export_figures_checkbox = widgets.Checkbox(value=True, description="Export Figures")

# Plot closest to mean
plot_closest_to_mean_checkbox = widgets.Checkbox(value=True, description="Use Closest-to-Mean Plotting")

# Export FOOOF results as pickle/Excel
export_fg_pickle_checkbox = widgets.Checkbox(value=True, description="Export as Pickle")
export_fg_excel_checkbox = widgets.Checkbox(value=False, description="Export as Excel")

# Display FOOOF results table
display_fg_export_table_checkbox = widgets.Checkbox(value=True, description="Display Export Table")

# Figure format checkboxes
fmt_options = ["png", "svg", "jpeg", "ppt"]
format_checkboxes = [widgets.Checkbox(value=False, description=fmt.upper()) for fmt in fmt_options]

# Buttons
file_picker_button = widgets.Button(description="Select Excel Files", button_style="info", icon="folder")
directory_picker_button = widgets.Button(description="Select Output Directory", button_style="info", icon="folder")
run_button = widgets.Button(description="Process and Analyze", button_style="success", icon="check")

output = widgets.Output()

###################################
# 1. FILE SELECTION CALLBACK
###################################
def on_file_picker_button_click(b):
    global selected_files
    root = Tk()
    root.withdraw()  # Hide the root window
    file_paths = filedialog.askopenfilenames(filetypes=[("Excel files", "*.xlsx")])
    selected_files = list(file_paths)
    if selected_files:
        print(f"Selected files: {selected_files}")
    else:
        print("No files selected.")

file_picker_button.on_click(on_file_picker_button_click)

###################################
# 2. DIRECTORY SELECTION CALLBACK
###################################
def on_directory_picker_button_click(b):
    global export_directory
    root = Tk()
    root.withdraw()  # Hide the root window
    directory_path = filedialog.askdirectory()
    if directory_path:
        export_directory = directory_path
        print(f"Selected output directory: {export_directory}")
    else:
        print("No directory selected.")

directory_picker_button.on_click(on_directory_picker_button_click)

######################################################
# 3. MAIN BUTTON: PROCESS, ANALYZE, & EXPORT
######################################################
def on_run_button_click(b):
    global fg_dict
    with output:
        output.clear_output()

        # 1) Check if files are selected.
        if not selected_files:
            print("Error: Please select Excel files to process.")
            return

        # 2) If exporting figures or results, ensure an export directory is selected.
        if (export_figures_checkbox.value or export_fg_pickle_checkbox.value or export_fg_excel_checkbox.value) and not export_directory:
            print("Error: Please select an output directory for exporting.")
            return

        # 3) Figure formats
        if export_figures_checkbox.value:
            chosen_formats = [cb.description.lower() for cb in format_checkboxes if cb.value]
            if not chosen_formats:
                print("Warning: No figure format selected. Proceeding without figure export.")
        else:
            chosen_formats = []  # Show plots but do not export

        # 4) Process the Excel files according to analysis mode
        print(f"Analysis mode: {analysis_mode.value}")
        
        if analysis_mode.value == "Grouped Mice":
            psd_data_dict_all = process_single_sheet_excel(selected_files)
        elif analysis_mode.value == "Individual Mice":
            psd_data_dict_all = process_eye_averages_excel(selected_files)
        elif analysis_mode.value == "Individual Mice Means":
            psd_data_dict_all = process_individual_mice_means_excel(selected_files)
        else:
            psd_data_dict_all = {}

        if not psd_data_dict_all:
            print("No data was processed. Please check your files and try again.")
            return

        # 5) Run FOOOF analysis
        freq_range = [freq_range_min.value, freq_range_max.value]
        peak_widths = [peak_width_min.value, peak_width_max.value]

        print("Running FOOOF analysis...")
        fg_dict = run_fooof_analysis(
            psd_data_dict_all,
            freq_range=freq_range,
            amp_threshold=amp_threshold.value,
            r2_threshold=r2_threshold.value,
            max_peaks=max_peaks.value,
            fitting_mode=fitting_mode.value,
            peak_width_limits=peak_widths
        )

        # 6) Plot the results
        print("\nPlotting FOOOF results...")
        for dict_key, fg_item in fg_dict.items():
            try:
                print(f"\nPlotting results for: {dict_key}")
                if plot_closest_to_mean_checkbox.value and len(fg_item) > 1:
                    plot_closest_to_mean(
                        fg_item, dict_key, export_directory, dict_key, chosen_formats,
                        x_min=x_axis_min.value,
                        x_max=x_axis_max.value,
                        y_min=y_axis_min.value,
                        y_max=y_axis_max.value,
                        show_grid=grid_checkbox.value,
                        x_tick_font_size=x_tick_font_size_widget.value,
                        include_r2=include_r2_checkbox.value,
                        include_peak_table=include_peak_table_checkbox.value
                    )
                else:
                    plot_all_psds(
                        fg_item, dict_key, export_directory, dict_key, chosen_formats,
                        x_min=x_axis_min.value,
                        x_max=x_axis_max.value,
                        y_min=y_axis_min.value,
                        y_max=y_axis_max.value,
                        show_grid=grid_checkbox.value,
                        x_tick_font_size=x_tick_font_size_widget.value,
                        include_r2=include_r2_checkbox.value,
                        include_peak_table=include_peak_table_checkbox.value
                    )
            except Exception as e:
                print(f"Error plotting {dict_key}: {e}")

        # 7) Export FOOOF results
        if export_fg_pickle_checkbox.value:
            pickle_filename = os.path.join(export_directory, "fooof_groups.pkl")
            save_fooof_group(fg_dict, filename=pickle_filename)

        if export_fg_excel_checkbox.value:
            excel_filename = os.path.join(export_directory, "fooof_results.xlsx")
            export_fg_dict_excel(fg_dict, excel_filename)

        # 8) Optionally display a summary table
        if display_fg_export_table_checkbox.value:
            display_fg_dict_table(fg_dict)

run_button.on_click(on_run_button_click)

##############################
# 4. DISPLAY THE GUI
##############################
display(
    widgets.VBox(
        [
            widgets.Label(value="1) Select Excel Files:"),
            file_picker_button,
            widgets.Label(value="2) Select Output Directory (for exporting):"),
            directory_picker_button,
            widgets.Label(value="3) Choose Analysis Mode (Grouped vs. Individual Mice, or Individual Mice Means):"),
            analysis_mode,
            widgets.Label(value="4) Configure FOOOF Parameters:"),
            widgets.HBox([freq_range_min, freq_range_max]),
            amp_threshold,
            r2_threshold,
            max_peaks,
            fitting_mode,
            widgets.HBox([peak_width_min, peak_width_max]),
            widgets.Label(value="(Optional) Specify Plot Axis Ranges:"),
            widgets.HBox([x_axis_min, x_axis_max]),
            widgets.HBox([y_axis_min, y_axis_max]),
            widgets.Label(value="(Optional) Grid, X-Tick Font, R² in Titles, and Peak Tables:"),
            widgets.HBox([grid_checkbox, x_tick_font_size_widget, include_r2_checkbox, include_peak_table_checkbox]),
            widgets.Label(value="5) Figure Export Options:"),
            export_figures_checkbox,
            plot_closest_to_mean_checkbox,
            widgets.Label(value="Choose Figure Export Formats:"),
            widgets.HBox(format_checkboxes),
            widgets.Label(value="6) Export FOOOF Results:"),
            widgets.HBox([export_fg_pickle_checkbox, export_fg_excel_checkbox]),
            widgets.Label(value="7) Display FOOOF Group Export in the Notebook:"),
            display_fg_export_table_checkbox,
            run_button,
            output,
        ]
    )
)
