import os
import pickle
import numpy as np
import matplotlib.pyplot as plt
import math

import ipywidgets as widgets
from ipyfilechooser import FileChooser
from IPython.display import display, clear_output

# Check for python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Warning: python-pptx not installed, PPT export will be disabled.")

###############################################################################
# 1) Helper Functions
###############################################################################

def exclude_traces(
    psd_array,
    freqs,
    low_band=(1, 3),
    low_band_threshold=3.0,  # multiple of mean in low band
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,  # multiple of mean in test bands
    test_band_count_threshold=None
):
    """
    Identifies which PSD traces (epochs) should be excluded based on two criteria:
      1) Low-band outliers (e.g., trace in 1-3 Hz band > threshold * mean_psd).
      2) Repeated suprathreshold events in test bands (e.g., >= 10x the mean).
    Returns (kept_indices, excluded_indices).
    """
    if test_band_count_threshold is None:
        test_band_count_threshold = len(test_bands) // 2

    mean_psd = np.mean(psd_array, axis=0)
    excluded_traces = []
    kept_traces = []

    # Indices for the specified low frequency band
    low_band_indices = np.where((freqs >= low_band[0]) & (freqs <= low_band[1]))[0]
    # Indices for each test band
    band_indices = [
        np.where((freqs >= band[0]) & (freqs <= band[1]))[0]
        for band in test_bands
    ]

    for i, trace in enumerate(psd_array):
        # 1) Check low-frequency extreme outlier
        if np.any(trace[low_band_indices] > low_band_threshold * mean_psd[low_band_indices]):
            excluded_traces.append(i)
            continue

        # 2) Check repeated suprathreshold events in test bands
        suprathreshold_count = 0
        for indices in band_indices:
            if np.any(trace[indices] > test_band_threshold * mean_psd[indices]):
                suprathreshold_count += 1

        if suprathreshold_count >= test_band_count_threshold:
            excluded_traces.append(i)
        else:
            kept_traces.append(i)

    return kept_traces, excluded_traces

def plot_psds_with_exclusion(
    ax,
    psd_array,
    freqs,
    kept_traces,
    excluded_traces,
    original_mean_psd,
    title,
    # Booleans
    show_kept=True,
    show_excluded=True,
    show_original_mean=True,
    show_new_mean=True,
    show_vertical_lines=True,
    # Colors
    color_kept="lightgray",
    color_excluded="red",
    color_old_mean="blue",
    color_new_mean="green",
    alpha_excluded=0.05,
    alpha_kept=0.7,
    # Fonts
    title_fontsize=8,
    axis_label_fontsize=8,
    legend_fontsize=8,
    tick_label_fontsize=8,
    max_title_length=40,
    # Vertical lines
    vertical_lines=None,
    vertical_line_color="black",
    vertical_line_style="--",
    vertical_line_alpha=0.6
):
    """
    Plots a single channel's PSD, highlighting kept vs excluded epochs, showing
    original mean vs new mean, etc.
    """
    # Truncate the title if needed
    if len(title) > max_title_length:
        title = title[:max_title_length] + "..."

    # 1) Kept traces
    if show_kept:
        for idx_i, idx in enumerate(kept_traces):
            label = "Kept Trace" if idx_i == 0 else None
            ax.plot(freqs, psd_array[idx], color=color_kept, alpha=alpha_kept, label=label)

    # 2) Original mean
    if show_original_mean:
        ax.plot(freqs, original_mean_psd, color=color_old_mean, linewidth=2, label="Original Mean")

    # 3) New mean (mean of only kept traces)
    if show_new_mean and len(kept_traces) > 0:
        new_mean_psd = np.mean(psd_array[kept_traces], axis=0)
        ax.plot(freqs, new_mean_psd, color=color_new_mean, linewidth=2, label="New Mean")

    # 4) Excluded traces
    if show_excluded:
        for idx_j, idx in enumerate(excluded_traces):
            label = "Excluded Trace" if idx_j == 0 else None
            ax.plot(freqs, psd_array[idx], color=color_excluded, alpha=alpha_excluded, zorder=10, label=label)

    # 5) Vertical lines
    if show_vertical_lines and vertical_lines is not None:
        for vfreq in vertical_lines:
            ax.axvline(
                vfreq,
                color=vertical_line_color,
                linestyle=vertical_line_style,
                alpha=vertical_line_alpha
            )

    # 6) Labels, title, legend
    ax.set_xlabel("Frequency (Hz)", fontsize=axis_label_fontsize)
    ax.set_ylabel("PSD (V²/Hz)", fontsize=axis_label_fontsize)
    ax.set_title(title, fontsize=title_fontsize)
    ax.legend(loc="upper right", fontsize=legend_fontsize)
    ax.tick_params(axis='both', labelsize=tick_label_fontsize)

def plot_psds_with_dropped_traces(
    psds_dict,
    rows_of_psds,
    # Exclusion
    low_band=(1,3),
    low_band_threshold=3.0,
    test_bands=[(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)],
    test_band_threshold=10.0,
    test_band_count_threshold=None,
    # Ranges
    x_min=None,
    x_max=None,
    y_min=None,
    y_max=None,
    # Booleans
    show_kept=True,
    show_excluded=True,
    show_original_mean=True,
    show_new_mean=True,
    show_vertical_lines=True,
    # Colors
    color_kept="lightgray",
    color_excluded="red",
    color_old_mean="blue",
    color_new_mean="green",
    # Grid layout
    num_cols=4,
    # Fonts
    title_fontsize=8,
    axis_label_fontsize=8,
    legend_fontsize=8,
    tick_label_fontsize=8,
    max_title_length=40,
    # Vertical lines
    vertical_lines=None,
    vertical_line_color="black",
    vertical_line_style="--",
    vertical_line_alpha=0.6
):
    """
    Plots PSD data for multiple channels in a grid, returning:
      - figures: list of figure objects
      - kept_indices_dict: dict mapping {channel: [kept epoch indices]}
    """
    figures = []
    kept_indices_dict = {}

    for row_idx, row in enumerate(rows_of_psds, start=1):
        valid_keys = [k for k in row if k in psds_dict]
        num_plots = len(valid_keys)
        if num_plots == 0:
            continue

        num_rows = math.ceil(num_plots / num_cols)
        fig, axes = plt.subplots(num_rows, num_cols, figsize=(num_cols * 4, num_rows * 4))
        plt.subplots_adjust(hspace=0.5)

        if num_rows * num_cols == 1:
            axes = [axes]
        else:
            axes = axes.flatten()

        for idx, key in enumerate(valid_keys):
            ax = axes[idx]
            psd_output = psds_dict[key]
            psd_data, freq_data = psd_output.get('psd'), psd_output.get('freqs')
            if psd_data is None or freq_data is None:
                ax.text(0.5, 0.5, f"No PSD data for {key}", ha='center', va='center')
                continue

            original_mean_psd = np.mean(psd_data, axis=0)

            # 1) Exclude traces
            kept_traces, excluded_traces = exclude_traces(
                psd_data,
                freq_data,
                low_band=low_band,
                low_band_threshold=low_band_threshold,
                test_bands=test_bands,
                test_band_threshold=test_band_threshold,
                test_band_count_threshold=test_band_count_threshold
            )
            kept_indices_dict[key] = kept_traces

            # 2) Plot
            plot_psds_with_exclusion(
                ax=ax,
                psd_array=psd_data,
                freqs=freq_data,
                kept_traces=kept_traces,
                excluded_traces=excluded_traces,
                original_mean_psd=original_mean_psd,
                title=key,
                show_kept=show_kept,
                show_excluded=show_excluded,
                show_original_mean=show_original_mean,
                show_new_mean=show_new_mean,
                show_vertical_lines=show_vertical_lines,
                color_kept=color_kept,
                color_excluded=color_excluded,
                color_old_mean=color_old_mean,
                color_new_mean=color_new_mean,
                alpha_excluded=0.05,
                alpha_kept=0.7,
                title_fontsize=title_fontsize,
                axis_label_fontsize=axis_label_fontsize,
                legend_fontsize=legend_fontsize,
                tick_label_fontsize=tick_label_fontsize,
                max_title_length=max_title_length,
                vertical_lines=vertical_lines,
                vertical_line_color=vertical_line_color,
                vertical_line_style=vertical_line_style,
                vertical_line_alpha=vertical_line_alpha
            )

            # x/y range
            if x_min is not None or x_max is not None:
                ax.set_xlim(x_min, x_max)
            if y_min is not None or y_max is not None:
                ax.set_ylim(y_min, y_max)

        # Turn off any unused subplots
        for ax in axes[num_plots:]:
            ax.axis('off')

        plt.tight_layout(rect=[0, 0.03, 1, 0.95])
        plt.suptitle(f"Row {row_idx}", fontsize=title_fontsize + 2)
        figures.append(fig)
        plt.show()

    return figures, kept_indices_dict

def group_keys_by_rows(psd_keys, row_size=4):
    """
    Utility to chunk the psd_keys list into sublists of length row_size.
    """
    rows = []
    for i in range(0, len(psd_keys), row_size):
        rows.append(psd_keys[i:i+row_size])
    return rows


###############################################################################
# 2) Main GUI
###############################################################################

def build_exportable_plot_psd_gui():
    """
    GUI that loads PSD pickle, selects channels, sets threshold for excluding PSD
    epochs, then plots and optionally exports a cleaned version of the PSD.
    This version uses ipyfilechooser for figure export and cleaned-PSD export.
    """
    # ~~~ UI Elements ~~~
    load_psd_button = widgets.Button(description='Load PSD Pickle', button_style='info')
    psd_file_chooser = FileChooser(os.getcwd(), title='Select PSD Pickle File', select_default=False)
    psd_file_chooser.show_only_files = True
    psd_file_chooser.filter_pattern = ['*.pkl']

    channels_dropdown = widgets.SelectMultiple(
        options=[],
        description='Select Channels:',
        layout=widgets.Layout(width='300px', height='200px')
    )

    # Plot toggles
    show_kept_cb = widgets.Checkbox(value=True, description='Show Kept Traces')
    show_excluded_cb = widgets.Checkbox(value=True, description='Show Excluded Traces')
    show_original_mean_cb = widgets.Checkbox(value=True, description='Show Original Mean')
    show_new_mean_cb = widgets.Checkbox(value=True, description='Show New Mean')
    show_vertical_lines_cb = widgets.Checkbox(value=True, description='Show Vertical Lines')

    # Vertical lines
    vertical_lines_text = widgets.Text(value="10,15", description='Vertical Lines (Hz):', layout=widgets.Layout(width='300px'))

    # Axis ranges
    x_min_widget = widgets.FloatText(value=None, description='X Min:', layout=widgets.Layout(width='150px'))
    x_max_widget = widgets.FloatText(value=None, description='X Max:', layout=widgets.Layout(width='150px'))
    y_min_widget = widgets.FloatText(value=None, description='Y Min:', layout=widgets.Layout(width='150px'))
    y_max_widget = widgets.FloatText(value=None, description='Y Max:', layout=widgets.Layout(width='150px'))

    # Font sizes
    title_fs_widget = widgets.IntText(value=8, description='Title FS:', layout=widgets.Layout(width='120px'))
    axis_fs_widget = widgets.IntText(value=8, description='Axis FS:', layout=widgets.Layout(width='120px'))
    legend_fs_widget = widgets.IntText(value=8, description='Legend FS:', layout=widgets.Layout(width='120px'))
    tick_fs_widget = widgets.IntText(value=8, description='Tick FS:', layout=widgets.Layout(width='120px'))
    max_title_len_widget = widgets.IntText(value=40, description='Max Title:', layout=widgets.Layout(width='120px'))

    # Colors
    color_kept_widget = widgets.ColorPicker(value='lightgray', description='Kept Trace:', layout=widgets.Layout(width='150px'))
    color_excl_widget = widgets.ColorPicker(value='red', description='Excluded Trace:', layout=widgets.Layout(width='150px'))
    color_oldmean_widget = widgets.ColorPicker(value='blue', description='Original Mean:', layout=widgets.Layout(width='150px'))
    color_newmean_widget = widgets.ColorPicker(value='green', description='New Mean:', layout=widgets.Layout(width='150px'))

    #  ~~~  Low/High Threshold Widgets  ~~~
    low_band_min_widget = widgets.FloatText(value=1.0, description='LowBandMin:', layout=widgets.Layout(width='150px'))
    low_band_max_widget = widgets.FloatText(value=3.0, description='LowBandMax:', layout=widgets.Layout(width='150px'))
    low_band_thr_widget = widgets.FloatText(value=3.0, description='LowBand Thr:', layout=widgets.Layout(width='150px'))
    test_band_thr_widget = widgets.FloatText(value=10.0, description='TestBand Thr:', layout=widgets.Layout(width='150px'))

    # Test bands
    test_band_text = widgets.Textarea(
        value="(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)",
        description='Test Bands:',
        layout=widgets.Layout(width='300px', height='80px')
    )

    num_cols_widget = widgets.IntText(value=4, description='Cols/Row:', layout=widgets.Layout(width='130px'))

    plot_psd_button = widgets.Button(description='Plot PSDs', button_style='success')

    # ~~~ Figure Export ~~~
    export_fig_chooser = FileChooser(
        os.getcwd(),
        title='Select or type a file name (e.g. MyFigure.png or MyPresentation.pptx)',
        select_default=False
    )
    # We can allow any extension, but you might prefer restricting:
    # export_fig_chooser.filter_pattern = ['*.png','*.jpeg','*.jpg','*.svg','*.pptx']

    export_button = widgets.Button(description='Export Figures', button_style='warning')

    # ~~~ Cleaned PSD Export ~~~
    cleaned_pickle_chooser = FileChooser(
        os.getcwd(),
        title='Select or type a .pkl file for cleaned data export',
        select_default=False
    )
    cleaned_pickle_chooser.filter_pattern = ['*.pkl']
    export_cleaned_button = widgets.Button(description='Export Cleaned PSD', button_style='info')

    # Output areas
    load_output_area = widgets.Output()
    plot_output_area = widgets.Output()
    export_output_area = widgets.Output()

    # Storage
    loaded_psd = {}
    current_figures = []
    kept_indices_dict = {}

    # 1) LOAD PSD
    def on_load_psd_clicked(b):
        with load_output_area:
            clear_output()
            chosen_file = psd_file_chooser.selected
            if not chosen_file:
                print("No file selected.")
                return
            if not os.path.isfile(chosen_file):
                print(f"Invalid file: {chosen_file}")
                return
            try:
                with open(chosen_file, 'rb') as f:
                    loaded_psd.clear()
                    loaded_psd.update(pickle.load(f))
                print(f"Loaded PSD from '{chosen_file}'")

                # If your channels are named "Ch1", "Ch2", etc., you can sort them numerically:
                # Otherwise, you can remove or adapt the sorting key as needed
                channels = sorted(loaded_psd.keys(), key=lambda x: int(x.replace("Ch","")))
                channels_dropdown.options = channels
                print(f"Channels: {channels}")
            except Exception as e:
                print(f"Error loading PSD: {e}")

    load_psd_button.on_click(on_load_psd_clicked)

    # 2) PLOT PSD
    def on_plot_psd_clicked(b):
        nonlocal current_figures, kept_indices_dict
        with plot_output_area:
            clear_output()

            selected_channels = list(channels_dropdown.value)
            if not selected_channels:
                print("No channels selected.")
                return

            # Booleans
            show_kept = show_kept_cb.value
            show_excluded = show_excluded_cb.value
            show_original_mean = show_original_mean_cb.value
            show_new_mean = show_new_mean_cb.value
            show_vertical_lines = show_vertical_lines_cb.value

            # Axis
            x_min = x_min_widget.value
            x_max = x_max_widget.value
            y_min = y_min_widget.value
            y_max = y_max_widget.value

            # Font
            title_fs = title_fs_widget.value
            axis_fs = axis_fs_widget.value
            legend_fs = legend_fs_widget.value
            tick_fs = tick_fs_widget.value
            max_title_length = max_title_len_widget.value

            # Colors
            color_kept = color_kept_widget.value
            color_excl = color_excl_widget.value
            color_old_mean = color_oldmean_widget.value
            color_new_mean = color_newmean_widget.value

            # Low band
            lb_min = low_band_min_widget.value
            lb_max = low_band_max_widget.value
            low_band_threshold = low_band_thr_widget.value

            # test band threshold
            tb_threshold = test_band_thr_widget.value

            # Test bands
            test_str = test_band_text.value.strip()
            test_bands_list = []
            if test_str:
                test_str = test_str.replace(" ", "")
                pairs = test_str.split(")")
                for p in pairs:
                    p = p.strip(",").strip().strip("(").strip()
                    if not p:
                        continue
                    vals = p.split(",")
                    if len(vals) == 2:
                        try:
                            lowf = float(vals[0])
                            highf = float(vals[1])
                            test_bands_list.append((lowf, highf))
                        except:
                            pass
                if not test_bands_list:
                    test_bands_list = [(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)]
            else:
                test_bands_list = [(7,9),(9,11),(11,13),(13,15),(15,17),(17,19),(19,21)]

            num_cols = num_cols_widget.value
            if num_cols < 1:
                num_cols = 4

            # Vertical lines
            vlines_str = vertical_lines_text.value.strip()
            vert_lines_list = []
            if vlines_str:
                parts = vlines_str.split(",")
                for part in parts:
                    part = part.strip()
                    if part:
                        try:
                            freq_val = float(part)
                            vert_lines_list.append(freq_val)
                        except:
                            pass

            # Build data dict for only selected channels
            psd_data_dict = {}
            for ch in selected_channels:
                ch_info = loaded_psd.get(ch, {})
                psd_array = ch_info.get('psd', None)
                freqs = ch_info.get('freqs', None)
                if psd_array is not None and freqs is not None:
                    psd_data_dict[ch] = {'psd': psd_array, 'freqs': freqs}

            # group into rows
            rows_of_psds = group_keys_by_rows(selected_channels, row_size=num_cols)

            # Call the function
            figs, kept_dict = plot_psds_with_dropped_traces(
                psds_dict=psd_data_dict,
                rows_of_psds=rows_of_psds,
                low_band=(lb_min, lb_max),
                low_band_threshold=low_band_threshold,
                test_bands=test_bands_list,
                test_band_threshold=tb_threshold,
                test_band_count_threshold=None,
                x_min=x_min,
                x_max=x_max,
                y_min=y_min,
                y_max=y_max,
                show_kept=show_kept,
                show_excluded=show_excluded,
                show_original_mean=show_original_mean,
                show_new_mean=show_new_mean,
                show_vertical_lines=show_vertical_lines,
                color_kept=color_kept,
                color_excluded=color_excl,
                color_old_mean=color_old_mean,
                color_new_mean=color_new_mean,
                num_cols=num_cols,
                title_fontsize=title_fs,
                axis_label_fontsize=axis_fs,
                legend_fontsize=legend_fs,
                tick_label_fontsize=tick_fs,
                max_title_length=max_title_length,
                vertical_lines=vert_lines_list,
                vertical_line_color="black",
                vertical_line_style="--",
                vertical_line_alpha=0.6
            )

            current_figures = figs
            kept_indices_dict = kept_dict

            if figs:
                print(f"Plotted {len(figs)} figure(s) for {len(selected_channels)} channel(s).")
                print("Kept indices recorded for each channel.")
            else:
                print("No figures plotted.")

    plot_psd_button.on_click(on_plot_psd_clicked)

    # 3) EXPORT FIGURES
    def on_export_button_clicked(b):
        with export_output_area:
            clear_output()
            if not current_figures:
                print("No figures to export (plot first).")
                return

            chosen_export_path = export_fig_chooser.selected
            if not chosen_export_path:
                print("No export path chosen. Please select or type a file name.")
                return

            # Deduce extension
            base, ext = os.path.splitext(chosen_export_path)
            ext = ext.lower()

            if not ext:
                print("No file extension detected (e.g., .png, .pptx). Please include one.")
                return

            # If PPTX => export slides
            if ext in ('.ppt', '.pptx'):
                if not HAS_PPTX:
                    print("python-pptx not installed. Cannot export PPT.")
                    return
                ppt_file = chosen_export_path
                print(f"Exporting {len(current_figures)} figure(s) to {ppt_file}")
                prs = Presentation()
                # Slide layout
                blank_layout = prs.slide_layouts[6]  # blank
                for idx, fig in enumerate(current_figures, start=1):
                    # Save temp PNG
                    temp_png = f"{base}_temp_{idx}.png"
                    fig.savefig(temp_png, format='png', dpi=150)
                    slide = prs.slides.add_slide(blank_layout)
                    left = top = Inches(1)
                    slide.shapes.add_picture(temp_png, left, top, width=Inches(8), height=Inches(4.5))
                    # remove temp
                    os.remove(temp_png)
                prs.save(ppt_file)
                print(f"PPTX exported => {ppt_file}")
            else:
                # Otherwise, treat as an image format (png, jpeg, svg, etc.)
                # We'll export each figure with a suffix
                print(f"Exporting {len(current_figures)} figure(s) as '{ext}' images.")
                for i, fig in enumerate(current_figures, start=1):
                    out_file = f"{base}_{i}{ext}"
                    fig.savefig(out_file, format=ext.lstrip('.'), dpi=150)
                    print(f"Saved => {out_file}")

            print("Done exporting figures.")

    export_button.on_click(on_export_button_clicked)

    # 4) EXPORT CLEANED PSD
    def on_export_cleaned_clicked(b):
        with export_output_area:
            clear_output()
            if not kept_indices_dict:
                print("No kept indices found. Please plot first.")
                return
            if not loaded_psd:
                print("No PSD data loaded.")
                return

            # Let user choose or type the .pkl for saving
            chosen_pickle_path = cleaned_pickle_chooser.selected
            if not chosen_pickle_path:
                print("No path for cleaned PSD chosen.")
                return

            if not chosen_pickle_path.lower().endswith('.pkl'):
                print("Please select or type a file ending in .pkl")
                return

            # Build the "cleaned_data" dict
            cleaned_data = {}
            for ch, ch_info in loaded_psd.items():
                if ch not in kept_indices_dict:
                    # Possibly skip or keep as original. We'll skip here
                    continue
                kept_ix = kept_indices_dict[ch]
                psd = ch_info.get('psd')
                freqs = ch_info.get('freqs')
                if psd is None or freqs is None:
                    continue
                # Keep only the kept epochs
                cleaned_psd = psd[kept_ix]
                cleaned_data[ch] = {
                    'psd': cleaned_psd,
                    'freqs': freqs
                }

            if not cleaned_data:
                print("No cleaned data generated.")
                return

            try:
                with open(chosen_pickle_path, 'wb') as f:
                    pickle.dump(cleaned_data, f)
                print(f"Cleaned PSD saved => {chosen_pickle_path}")
                print(f"Channels in cleaned data: {list(cleaned_data.keys())}")
            except Exception as e:
                print(f"Error saving cleaned PSD: {e}")

    export_cleaned_button.on_click(on_export_cleaned_clicked)

    # ~~~~~~~~~~~ Layout ~~~~~~~~~~~
    load_box = widgets.VBox([
        widgets.HTML("<h3>Load PSD</h3>"),
        load_psd_button,
        psd_file_chooser,
        load_output_area
    ])

    channel_sel_box = widgets.VBox([
        widgets.HTML("<h3>Select Channels</h3>"),
        channels_dropdown
    ])

    # Add Exclusion Parameter widgets
    exclusion_params_box = widgets.HBox([
        low_band_min_widget,
        low_band_max_widget,
        low_band_thr_widget,
        test_band_thr_widget
    ])

    # Collect controls for plotting
    plot_ctrls_box = widgets.VBox([
        widgets.HTML("<h3>Plotting Options</h3>"),
        show_kept_cb,
        show_excluded_cb,
        show_original_mean_cb,
        show_new_mean_cb,
        show_vertical_lines_cb,
        vertical_lines_text,

        widgets.HTML("<b>Exclusion Params</b>"),
        exclusion_params_box,

        widgets.HTML("<b>Axis Ranges</b>"),
        widgets.HBox([x_min_widget, x_max_widget, y_min_widget, y_max_widget]),

        widgets.HTML("<b>Font Sizes</b>"),
        widgets.HBox([title_fs_widget, axis_fs_widget, legend_fs_widget, tick_fs_widget, max_title_len_widget]),

        widgets.HTML("<b>Colors</b>"),
        widgets.HBox([color_kept_widget, color_excl_widget, color_oldmean_widget, color_newmean_widget]),

        widgets.HTML("<b>Test Bands</b>"),
        test_band_text,

        widgets.HTML("<b>Columns per Row</b>"),
        num_cols_widget,

        plot_psd_button,
        plot_output_area
    ])

    export_box = widgets.VBox([
        widgets.HTML("<h3>Export Figures</h3>"),
        export_fig_chooser,
        export_button,
        widgets.HTML("<h3>Export Cleaned PSD</h3>"),
        cleaned_pickle_chooser,
        export_cleaned_button,
        export_output_area
    ])

    ui = widgets.VBox([
        widgets.HTML("<h2>PSD Plotter GUI (Threshold-based Exclusion + Cleaned Export)</h2>"),
        widgets.HBox([load_box, channel_sel_box]),
        plot_ctrls_box,
        export_box
    ])

    display(ui)

# Run the GUI
build_exportable_plot_psd_gui()
